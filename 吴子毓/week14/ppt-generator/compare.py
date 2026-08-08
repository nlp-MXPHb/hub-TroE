"""v1 vs v2 对比：代码 token、LLM token、生成耗时、PPT 质量。
用法：python compare.py
依赖：tiktoken, python-pptx, openai + 环境变量 DASHSCOPE_API_KEY
"""
import os, sys, time, json, importlib.util
os.environ["NO_PROXY"] = "*"
os.environ["no_proxy"] = "*"

try:
    import tiktoken
    enc = tiktoken.get_encoding("cl100k_base")
    def token_count(text):
        return len(enc.encode(text))
except Exception:
    def token_count(text):
        return len(text) // 3  # 粗略近似

from pptx import Presentation
from pptx.util import Emu

BASE = os.path.dirname(os.path.abspath(__file__))
V1 = os.path.join(BASE, "_initial_version")
V2 = BASE
TEMPLATE = os.path.join(BASE, "..", "..", "..", "自进化agent.pptx")
TOPIC = "自进化 Agent：从对话失败中自动演化的 Skill 机制"
FILES = ["SKILL.md", "generate_ppt.py", "outline_generator.py"]


def code_metrics(path):
    txt = open(path, encoding="utf-8").read()
    return {"lines": txt.count("\n") + 1, "chars": len(txt), "tokens": token_count(txt)}


def load_module(path, name):
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def verify_pptx(path):
    p = Presentation(path)
    s = p.slides[0]
    colors, fonts = set(), set()
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    try:
                        colors.add(str(r.font.color.rgb))
                    except Exception:
                        pass
                    if r.font.name:
                        fonts.add(r.font.name)
    return {
        "pages": len(p.slides),
        "size": f"{round(Emu(p.slide_width).inches, 2)}x{round(Emu(p.slide_height).inches, 2)}",
        "colors": sorted(colors),
        "fonts": sorted(fonts),
    }


# === 一、代码 Token 对比 ===
print("=" * 60)
print("一、代码 Token 对比（tiktoken cl100k_base）")
print("=" * 60)
code_res = {}
for label, d in [("v1 初始版", V1), ("v2 优化版", V2)]:
    files = {fn: code_metrics(os.path.join(d, fn)) for fn in FILES}
    total = sum(f["tokens"] for f in files.values())
    code_res[label] = {"files": files, "total": total}
    print(f"\n[{label}]")
    for fn, m in files.items():
        print(f"  {fn:<22} lines={m['lines']:>4}  chars={m['chars']:>5}  tokens={m['tokens']:>5}")
    print(f"  {'合计':<22}                                tokens={total:>5}")

# === 二、LLM 调用 Token 对比 ===
print("\n" + "=" * 60)
print("二、LLM 调用 Token 对比（实际调用 dashscope）")
print("=" * 60)
llm_res = {}
for label, d in [("v1", V1), ("v2", V2)]:
    mod = load_module(os.path.join(d, "outline_generator.py"), f"og_{label}")
    t0 = time.perf_counter()
    outline, usage = mod.generate_outline(TOPIC, 5)
    elapsed = (time.perf_counter() - t0) * 1000
    llm_res[label] = {
        "outline": outline, "prompt": usage.prompt_tokens,
        "completion": usage.completion_tokens, "total": usage.total_tokens,
        "elapsed_ms": elapsed,
    }
    print(f"\n[{label}]")
    print(f"  prompt_tokens     : {usage.prompt_tokens}")
    print(f"  completion_tokens : {usage.completion_tokens}")
    print(f"  total_tokens      : {usage.total_tokens}")
    print(f"  LLM 耗时          : {elapsed:.0f}ms")

# === 三、排版耗时与 PPT 质量 ===
print("\n" + "=" * 60)
print("三、排版耗时与 PPT 质量验证（动态模板，各跑 3 次取平均）")
print("=" * 60)
ppt_res = {}
for label, d in [("v1", V1), ("v2", V2)]:
    gen_mod = load_module(os.path.join(d, "generate_ppt.py"), f"gp_{label}")
    outline = llm_res[label]["outline"]
    out = os.path.join(BASE, f"compare_{label}.pptx")
    dts = []
    for _ in range(3):
        t0 = time.perf_counter()
        gen_mod.generate_ppt(outline, out, template_path=TEMPLATE)
        dts.append((time.perf_counter() - t0) * 1000)
    avg_ms = sum(dts) / len(dts)
    info = verify_pptx(out)
    ppt_res[label] = {"avg_ms": avg_ms, "info": info}
    print(f"\n[{label}]")
    print(f"  排版耗时 : avg={avg_ms:.1f}ms  (3次: {[f'{x:.0f}' for x in dts]})")
    print(f"  页数     : {info['pages']}")
    print(f"  尺寸     : {info['size']}")
    print(f"  封面配色 : {info['colors']}")
    print(f"  字体     : {info['fonts']}")

# === 四、优化幅度汇总 ===
print("\n" + "=" * 60)
print("四、优化幅度汇总")
print("=" * 60)


def pct(old, new):
    diff = 100 * (new - old) / old
    arrow = "降" if new < old else "升"
    return f"{arrow} {abs(diff):.1f}%"


v1c, v2c = code_res["v1 初始版"]["total"], code_res["v2 优化版"]["total"]
v1p, v2p = llm_res["v1"]["prompt"], llm_res["v2"]["prompt"]
v1cm, v2cm = llm_res["v1"]["completion"], llm_res["v2"]["completion"]
v1t, v2t = llm_res["v1"]["total"], llm_res["v2"]["total"]
v1e, v2e = ppt_res["v1"]["avg_ms"], ppt_res["v2"]["avg_ms"]
print(f"  代码 token 合计  : {v1c} -> {v2c}  ({pct(v1c, v2c)})")
print(f"  LLM prompt token : {v1p} -> {v2p}  ({pct(v1p, v2p)})")
print(f"  LLM completion   : {v1cm} -> {v2cm}  ({pct(v1cm, v2cm)})")
print(f"  LLM total token  : {v1t} -> {v2t}  ({pct(v1t, v2t)})")
print(f"  排版耗时         : {v1e:.1f}ms -> {v2e:.1f}ms  ({pct(v1e, v2e)})")

# 保存结果到 JSON
result = {
    "code": code_res, "llm": {k: {kk: vv for kk, vv in v.items() if kk != "outline"}
                              for k, v in llm_res.items()},
    "pptx": ppt_res,
}
with open(os.path.join(BASE, "compare_result.json"), "w", encoding="utf-8") as f:
    json.dump(result, f, ensure_ascii=False, indent=2, default=str)
print(f"\n✓ 详细结果已保存到 compare_result.json")
