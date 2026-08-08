"""PPT 生成器：大纲 → PPT，支持动态读取任意模板的配色/字体/尺寸规范。
依赖：python-pptx, lxml。传入 template_path 则自动套用该模板视觉规范。
大纲 layout: default | two_column | comparison | stack
"""
import json, sys, zipfile
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# 默认规范（无模板时 fallback）
DEFAULT_SPEC = {
    "w_in": 10.0, "h_in": 5.625,
    "font": "Calibri",
    "bg": "1A1A2E", "title": "FFFFFF", "accent": "2E86AB",
    "light": "99C4DD", "gray": "BBBBBB", "footer": "FFAAAA",
    "footer_text": "八斗学院出品|盗版必究",
}


def _hex(s):
    """'2E86AB' → RGBColor。"""
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def extract_template_spec(template_path):
    """从模板 pptx 提取视觉规范（尺寸/字体/配色/页脚），返回 spec dict。"""
    from collections import Counter
    tp = Presentation(template_path)
    spec = {
        "w_in": round(Emu(tp.slide_width).inches, 3),
        "h_in": round(Emu(tp.slide_height).inches, 3),
    }
    # 字体：从 theme 取主标题字体
    with zipfile.ZipFile(template_path) as z:
        root = etree.fromstring(z.read("ppt/theme/theme1.xml"))
    ns = {"a": A_NS}
    maj = root.find(".//a:fontScheme/a:majorFont/a:latin", ns)
    spec["font"] = maj.get("typeface") if maj is not None and maj.get("typeface") else "Calibri"
    # 扫描实际文本框，按字号范围 + 位置推断颜色角色
    title_colors, subtitle_colors, accent_colors, point_colors = [], [], [], []
    footer_colors = Counter()
    footer_text = ""
    for slide in list(tp.slides)[:5]:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if not t:
                continue
            is_footer = len(t) < 40 and ("出品" in t or "版权" in t or "©" in t or "盗版" in t)
            if is_footer:
                footer_text = t.replace("\n", "|")
            top_in = Emu(sh.top).inches if sh.top is not None else 0
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    try:
                        c = str(r.font.color.rgb)
                    except Exception:
                        continue
                    sz = r.font.size
                    if sz is None:
                        continue
                    sz_pt = int(sz.pt)
                    if is_footer:
                        footer_colors[c] += 1
                    elif sz_pt >= 30:
                        title_colors.append(c)
                    elif sz_pt >= 18:
                        subtitle_colors.append(c)
                    elif sz_pt >= 14 and top_in < 1.5:
                        accent_colors.append(c)    # 顶部小字→章节标识
                    else:
                        point_colors.append(c)     # 其余→要点

    def mode(lst, default):
        return Counter(lst).most_common(1)[0][0] if lst else default

    main_color = mode(title_colors + subtitle_colors + point_colors, None)
    spec["light"]  = mode(subtitle_colors, DEFAULT_SPEC["light"])
    # accent: 14pt 顶部排除主文字色（主文字常混入该位置），取剩余最常见
    accent_candidates = [c for c in accent_colors if c != main_color]
    spec["accent"] = mode(accent_candidates, DEFAULT_SPEC["accent"]) if accent_candidates else mode(accent_colors, DEFAULT_SPEC["accent"])
    spec["gray"]   = mode(point_colors, DEFAULT_SPEC["gray"])
    spec["footer"] = footer_colors.most_common(1)[0][0] if footer_colors else DEFAULT_SPEC["footer"]
    spec["footer_text"] = footer_text or DEFAULT_SPEC["footer_text"]
    # 背景：读显式背景；读不到则按主文字亮度推断（深文字→白底，浅文字→深底）
    spec["bg"] = DEFAULT_SPEC["bg"]
    try:
        bg = tp.slides[0].background
        if bg.fill.type is not None:
            spec["bg"] = str(bg.fill.fore_color.rgb)
    except Exception:
        pass
    if spec["bg"] == DEFAULT_SPEC["bg"] and main_color:
        ml = 0.299 * int(main_color[0:2], 16) + 0.587 * int(main_color[2:4], 16) + 0.114 * int(main_color[4:6], 16)
        spec["bg"] = "FFFFFF" if ml < 128 else DEFAULT_SPEC["bg"]
    # title: 保证与背景对比（浅底→用强调色，深底→用白）
    bg_lum = 0.299 * int(spec["bg"][0:2], 16) + 0.587 * int(spec["bg"][2:4], 16) + 0.114 * int(spec["bg"][4:6], 16)
    spec["title"] = spec["accent"] if bg_lum > 128 else mode(title_colors, DEFAULT_SPEC["title"])
    return spec


def _box(slide, spec, l, t, w, h, text, sz, color_key, bold=False):
    """放文本框，颜色按 color_key 从 spec 取。"""
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = spec["font"]
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = _hex(spec[color_key])


def _panel_box(slide, spec, l, t, w, h, fill_key=None, line_key=None, line_w=1.0):
    """分框底色矩形。fill_key/line_key 为 spec 色键，None 表示透明。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill_key and fill_key in spec:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex(spec[fill_key])
    else:
        shp.fill.background()
    if line_key and line_key in spec:
        shp.line.color.rgb = _hex(spec[line_key])
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp


def _luminance(hex_color):
    """hex 颜色亮度（0-255）。"""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _text_on(bg_hex):
    """根据底色亮度返回对比文字色键：浅底→深字，深底→浅字。
    中蓝/中灰这种中间亮度，用 gray（中灰，对比度稳定）。
    """
    lum = _luminance(bg_hex)
    if lum < 100:
        return "title"
    if lum < 180:
        return "gray"
    return "bg"


def _safe_text_rgb(bg_hex):
    """直接算一个绝对高对比的 RGB，绕过色键映射。
    底色 lum<128 → 写近白 (F0F0F0)；否则写近黑 (1A1A1A)。"""
    if _luminance(bg_hex) < 128:
        return RGBColor(0xF0, 0xF0, 0xF0)
    return RGBColor(0x1A, 0x1A, 0x1A)


def _box_safe(slide, spec, l, t, w, h, text, sz, bg_hex, bold=False):
    """_box 的安全版：文字色按 bg_hex 实际亮度算绝对对比色，不走色键。"""
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = spec["font"]
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = _safe_text_rgb(bg_hex)


def _render_header(slide, spec, item):
    """统一的头部（section / title / subtitle），各布局复用。"""
    _box(slide, spec, 0.5, 0.3, 9.0, 0.4, item.get("section", ""), 14, "accent", True)
    _box(slide, spec, 0.5, 0.8, 9.0, 0.8, item.get("title", ""), 32, "title", True)
    if item.get("subtitle"):
        _box(slide, spec, 0.5, 1.7, 9.0, 0.5, item["subtitle"], 18, "light")


def _render_default(slide, spec, item):
    """默认布局：原版"要点列表"渲染。"""
    _render_header(slide, spec, item)
    y, step, ph = 2.5, 0.55, 0.5
    for pt in item.get("points", []):
        _box(slide, spec, 0.7, y, 8.6, ph, "· " + pt, 16, "gray")
        y += step


def _render_two_column(slide, spec, item):
    """左右两栏：顶部标题区 + 副标题横条 + 两分框 + 底部高亮条。
    panel 字段：title/subtitle/body(body 可为字符串多行 或 [{label,color,fill}])。
    """
    _render_header(slide, spec, item)
    # 副标题横条
    if item.get("banner"):
        bg = spec.get("light", spec["bg"])
        _panel_box(slide, spec, 0.5, 2.4, 9.0, 0.5, fill_key="light", line_key=None)
        _box_safe(slide, spec, 0.7, 2.45, 8.6, 0.4, item["banner"], 16, bg, True)
    panels = item.get("panels", [])
    if not panels:
        return _render_default(slide, spec, item)
    col_w = 4.3
    gap = 0.4
    y0 = 3.0
    h_panel = 2.2
    for i, panel in enumerate(panels[:2]):
        l = 0.5 + i * (col_w + gap)
        _panel_box(slide, spec, l, y0, col_w, h_panel,
                   fill_key=panel.get("fill"),
                   line_key=panel.get("border", "gray"), line_w=1.5)
        # panel 实际底色
        panel_bg = spec.get(panel.get("fill", ""), spec["bg"]) if panel.get("fill") else spec["bg"]
        _box_safe(slide, spec, l + 0.2, y0 + 0.1, col_w - 0.4, 0.4,
                  panel.get("title", ""), 14, panel_bg, True)
        if panel.get("subtitle"):
            _box_safe(slide, spec, l + 0.2, y0 + 0.55, col_w - 0.4, 0.4,
                      panel["subtitle"], 16, panel_bg)
        body = panel.get("body", "")
        if isinstance(body, str):
            for j, line in enumerate(body.split("\n")):
                _box_safe(slide, spec, l + 0.2, y0 + 1.05 + j * 0.22,
                          col_w - 0.4, 0.22, line, 12, panel_bg)
        elif isinstance(body, list):
            for j, lv in enumerate(body):
                color = lv.get("color") or lv.get("fill") or "gray"
                _panel_box(slide, spec, l + 0.2, y0 + 1.05 + j * 0.24,
                           0.18, 0.20, fill_key=color, line_key=None)
                lv_bg = spec.get(color, spec["bg"])
                _box_safe(slide, spec, l + 0.45, y0 + 1.05 + j * 0.24,
                          col_w - 0.65, 0.22, lv.get("label", ""), 12, lv_bg, True)
    if item.get("footer"):
        _panel_box(slide, spec, 0.5, 5.3, 9.0, 0.32, fill_key="accent", line_key=None)
        _box_safe(slide, spec, 0.7, 5.32, 8.6, 0.3, item["footer"], 12, spec["accent"], True)


def _render_comparison(slide, spec, item):
    """左右对比：同一对象的两个版本（v1 vs v2）。
    item.left / item.right = {"title": str, "rows": [{"label": str, "value": str}, ...]}
    行高 0.45，逐行左半 + 右半对齐排版。
    """
    _render_header(slide, spec, item)
    left = item.get("left", {})
    right = item.get("right", {})
    if not left or not right:
        return _render_default(slide, spec, item)
    col_w = 4.3
    gap = 0.4
    y0 = 2.4
    # 栏头 accent 底 → 用绝对对比色（白/黑）
    for i, side in enumerate([left, right]):
        l = 0.5 + i * (col_w + gap)
        _panel_box(slide, spec, l, y0, col_w, 0.5, fill_key="accent", line_key=None)
        _box_safe(slide, spec, l + 0.2, y0 + 0.05, col_w - 0.4, 0.4,
                  side.get("title", ""), 16, spec["accent"], True)
    # 行（背景是白底 bg FFFFFF，文字用 light/title 也看得清）
    rows = max(len(left.get("rows", [])), len(right.get("rows", [])))
    row_h = 0.45
    y = y0 + 0.6
    for r in range(rows):
        lr = left.get("rows", [{}])[r] if r < len(left.get("rows", [])) else {}
        rr = right.get("rows", [{}])[r] if r < len(right.get("rows", [])) else {}
        for i, row in enumerate([lr, rr]):
            l = 0.5 + i * (col_w + gap)
            _box(slide, spec, l + 0.2, y, col_w * 0.55 - 0.2, row_h,
                 row.get("label", ""), 13, "light")
            _box(slide, spec, l + col_w * 0.55, y, col_w * 0.45 - 0.1, row_h,
                 row.get("value", ""), 13, "title", True)
        y += row_h
    # 底部结论（light 浅底）
    if item.get("footer"):
        _panel_box(slide, spec, 0.5, 5.3, 9.0, 0.32, fill_key="light", line_key=None)
        _box_safe(slide, spec, 0.7, 5.32, 8.6, 0.3, item["footer"], 12, spec["light"], True)


def _render_stack(slide, spec, item):
    """层栈/漏斗/金字塔。
    item.layers = [{"label": str, "level": int(1-5), "color": "色键"}]
    item.shape = "funnel" | "pyramid" | "flat"（默认 funnel：高层=宽）
    每层宽度按归一化 level 算，居中对称。
    """
    _render_header(slide, spec, item)
    layers = item.get("layers", [])
    if not layers:
        return _render_default(slide, spec, item)
    shape = item.get("shape", "funnel")
    n = len(layers)
    y0 = 2.4
    h_layer = 0.42
    gap = 0.08
    max_w = 8.6  # 顶层宽度
    min_w = 3.6  # 底层宽度
    center = 5.0
    for i, lv in enumerate(layers):
        level = lv.get("level", 1)
        if shape == "pyramid":
            ratio = level / max(l["level"] for l in layers)
            w = min_w + (max_w - min_w) * ratio
        elif shape == "flat":
            w = max_w
        else:
            ratio = level / max(l["level"] for l in layers)
            w = min_w + (max_w - min_w) * ratio
        l = center - w / 2
        t = y0 + i * (h_layer + gap)
        color = lv.get("color", "accent")
        _panel_box(slide, spec, l, t, w, h_layer, fill_key=color, line_key=None)
        # 文字色按该层实际底色取绝对对比色
        layer_bg = spec.get(color, spec["bg"])
        _box_safe(slide, spec, l + 0.2, t + 0.04, w - 0.4, h_layer - 0.05,
                  lv.get("label", ""), 13, layer_bg, True)
    if item.get("footer"):
        _panel_box(slide, spec, 0.5, 5.3, 9.0, 0.32, fill_key="accent", line_key=None)
        _box_safe(slide, spec, 0.7, 5.32, 8.6, 0.3, item["footer"], 12, spec["accent"], True)


LAYOUT_RENDERERS = {
    "default": _render_default,
    "two_column": _render_two_column,
    "comparison": _render_comparison,
    "stack": _render_stack,
}


def add_slide(prs, spec, item, is_cover=False):
    """加一页。is_cover=True 为封面；其余按 item.get('layout') 选渲染器。"""
    s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _hex(spec["bg"])
    if is_cover:
        _box(s, spec, 0.5, 1.5, 9.0, 1.2, item.get("title", ""), 44, "title", True)
        _box(s, spec, 0.5, 2.8, 9.0, 0.6, item.get("subtitle", ""), 22, "light")
        y, step, ph = 3.6, 0.45, 0.4
        for pt in item.get("points", []):
            _box(s, spec, 0.7, y, 8.6, ph, "· " + pt, 16, "gray")
            y += step
    else:
        layout = item.get("layout", "default")
        renderer = LAYOUT_RENDERERS.get(layout, _render_default)
        renderer(s, spec, item)
    _box(s, spec, 7.85, 0.15, 2.0, 0.4, spec["footer_text"], 12, "footer", True)
    return s


def generate_ppt(outline, output_path, template_path=None):
    """根据大纲生成 PPT。template_path 给定则动态套用该模板规范。"""
    spec = extract_template_spec(template_path) if template_path else dict(DEFAULT_SPEC)
    prs = Presentation()
    prs.slide_width = Inches(spec["w_in"])
    prs.slide_height = Inches(spec["h_in"])
    for i, item in enumerate(outline):
        add_slide(prs, spec, item, is_cover=(i == 0))
    prs.save(output_path)
    return output_path


SAMPLE_OUTLINE = [
    {"section": "Nous Research", "title": "自进化 Agent", "subtitle": "从对话失败中自动演化的 Skill 机制",
     "points": ["22% → 80% 准确率提升", "8 块 Nudge 进化轨迹", "契约式评估设计"]},
    {"section": "01 核心机制", "title": "Skill 是可进化的操作 SOP",
     "subtitle": "不是 RAG 知识搬运，而是带边界条件的决策流程",
     "points": ["Skill = 结构化操作文档，按步骤照做而非推理", "三种演化：patch 补分支 / create 建新域 / patch 重写强度", "Reviewer 仅接收失败样本，最小改动原则"]},
    {"section": "02 评估契约", "title": "Agent-Evaluator 契约式设计",
     "subtitle": "三类互斥失败原因，评估简单可解释",
     "points": ["不能答时只说「需要联系人工客服」", "评估器：推脱一票否决 / 缺关键词 / 出现禁止词", "数字归一化 + 否定前置检测，确定性可复现"]},
    {"section": "03 进化效果", "title": "准确率从 22% 提升到 80%",
     "subtitle": "每块失败驱动对应 Skill 出现",
     "points": ["vip_refund：25% → 100%（patch VIP 分支）", "digital_goods：0% → 100%（create 新 Skill）", "logistics/payment：0% → 60-100%（create 新 Skill）"]},
]


if __name__ == "__main__":
    args = sys.argv[1:]
    template = None
    if "--template" in args:
        i = args.index("--template")
        template = args[i + 1]
        args = args[:i] + args[i + 2:]
    if args:
        outline = json.loads(Path(args[0]).read_text(encoding="utf-8"))
        out = args[1] if len(args) > 1 else "output.pptx"
    else:
        outline, out = SAMPLE_OUTLINE, "sample_output.pptx"
    generate_ppt(outline, out, template_path=template)
    src = f"模板={template}" if template else "默认规范"
    print(f"已生成 PPT：{out}（共 {len(outline)} 页，{src}）")
